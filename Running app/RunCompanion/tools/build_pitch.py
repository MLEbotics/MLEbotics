from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()

# Slide 1 Cover
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Running Companion Project"
slide.placeholders[1].text = "Hybrid Native + Cloud Robotics Platform\nFebruary 27, 2026"

# Slide 2 Agenda
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Agenda"
body = slide.shapes.placeholders[1].text_frame
body.text = "Vision + Problem Statement"
for item in [
    "Solution overview & hybrid architecture",
    "Voice + public alerting differentiators",
    "Hardware & build plan",
    "Collaboration and funding asks",
    "Next steps",
]:
    body.add_paragraph().text = item

# Slide 3 Problem
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Problem Space"
body = slide.shapes.placeholders[1].text_frame
body.text = (
    "Outdoor runners lack live pacing support, safety signaling, and adaptive coaching while training alone."
)
for item in [
    "Inconsistent pacing and training efficiency",
    "Crowded routes with no proactive runner-approaching alerts",
    "Hands busy—no easy way to control devices mid-run",
    "Wearables collect data but cannot guide physical pace",
]:
    body.add_paragraph().text = item

# Slide 4 Solution
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Solution Overview"
body = slide.shapes.placeholders[1].text_frame
body.text = (
    "RunBot: a 4-wheel ESP32 robot that physically leads the runner while a Flutter app handles pace control, AI coaching, and remote alerts."
)
for item in [
    "Robot follows GPS waypoints at target speed with adaptive slowdown",
    "Owner-only voice commands (stop, speed up, alert ahead, return)",
    "Garmin BLE telemetry feeds live HR, cadence, distance",
    "AI coach + SMS broadcasting keep people informed ahead on the route",
]:
    body.add_paragraph().text = item

# Slide 5 Architecture diagram
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Hybrid Architecture"
left, top, width, height = Inches(0.5), Inches(1.6), Inches(3.8), Inches(1.6)
shape_native = slide.shapes.add_shape(1, left, top, width, height)
shape_native.fill.solid()
shape_native.fill.fore_color.rgb = RGBColor(0, 122, 116)
shape_native.line.color.rgb = RGBColor(255, 255, 255)
shape_native.text = "Native Core (Flutter)\n• Robot control\n• BLE & sensors\n• Voice + alerts"
shape_native.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
shape_native.text_frame.paragraphs[0].font.size = Pt(16)

shape_cloud = slide.shapes.add_shape(1, Inches(5.2), Inches(1.6), Inches(3.8), Inches(1.6))
shape_cloud.fill.solid()
shape_cloud.fill.fore_color.rgb = RGBColor(15, 118, 110)
shape_cloud.line.color.rgb = RGBColor(255, 255, 255)
shape_cloud.text = "Cloud Modules CDN\n• Investor decks\n• Onboarding flows\n• Per-partner content"
shape_cloud.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
shape_cloud.text_frame.paragraphs[0].font.size = Pt(16)

shape_robot = slide.shapes.add_shape(1, Inches(3.1), Inches(3.5), Inches(2.8), Inches(1.2))
shape_robot.fill.solid()
shape_robot.fill.fore_color.rgb = RGBColor(245, 130, 32)
shape_robot.line.color.rgb = RGBColor(255, 255, 255)
shape_robot.text = "RunBot Hardware\nESP32 + sensors"
shape_robot.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
shape_robot.text_frame.paragraphs[0].font.size = Pt(16)

slide.shapes.add_connector(1, Inches(2.2), Inches(3.2), Inches(3.3), Inches(3.5)).line.color.rgb = RGBColor(80, 80, 80)
slide.shapes.add_connector(1, Inches(7.0), Inches(3.2), Inches(4.7), Inches(3.5)).line.color.rgb = RGBColor(80, 80, 80)
textbox = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(8.0), Inches(0.6))
textbox.text_frame.text = (
    "Cloud manifest updates UI instantly; native core handles safety-critical operations."
)

# Slide 6 Hybrid frameworks explanation
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Hybrid Frameworks Loading Cloud Modules"
body = slide.shapes.placeholders[1].text_frame
body.text = "Middle ground approach:" 
for item in [
    "App built with Flutter/React Native/Ionic",
    "Most UI/logic pulled from our server via manifest",
    "Updates roll out instantly without app-store approval",
]:
    body.add_paragraph().text = item
body.add_paragraph().text = "Pros: better performance than WebView, supports cloud updates, app stores accept it."
body.add_paragraph().text = "Cons: more complex architecture, requires secure module-loading system."
body.add_paragraph().text = "Examples: Tesla app, older Airbnb builds, many enterprise field apps."

# Slide 7 Voice & Alerting
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Voice + Public Alerting"
body = slide.shapes.placeholders[1].text_frame
body.text = "Owner-only speech commands trigger pacing, alerts, and music without touching the phone."
for item in [
    "Voice command → robot speaker announcement → SMS broadcast to contacts ahead",
    "Directional alerts (on your left/right) to clear crowd space",
    "AI coach provides pacing feedback, obstacle warnings, and encouragement",
]:
    body.add_paragraph().text = item

# Slide 8 Hardware Roadmap
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Hardware Build Plan"
textbox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.4), Inches(1.0))
tf = textbox.text_frame
tf.text = "Bill of Materials (~$75):"
for item in [
    "Freenove 4WD ESP32 kit, NEO-6M GPS, SG90 servo, ultrasonic sensors",
    "Quality 18650 cells + wiring harness, optional BT speaker upgrade",
    "Firmware ready: GPS pacing, obstacle avoidance, runner distance monitoring",
]:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

timeline = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(8.4), Inches(2.0))
tf = timeline.text_frame
tf.text = "Execution timeline:"
for item in [
    "Week 1 — Procurement + assembly",
    "Week 2 — Firmware flash + integration tests",
    "Week 3 — Outdoor validation (pacing + alerting + voice)",
]:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

# Slide 9 Hardware Sourcing Strategy
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Hardware Sourcing Strategy"
body = slide.shapes.placeholders[1].text_frame
body.text = "Blend off-the-shelf robotics kits with custom pacer chassis to accelerate pilots."
for item in [
    "Phase 1: buy Freenove/DFRobot components locally for rapid prototyping",
    "Phase 2: standardize BOM and negotiate bulk pricing with Shenzhen suppliers",
    "Phase 3: dedicated pacer chassis + mobile aid station payload (water, gels, med kit)",
]:
    body.add_paragraph().text = item
body.add_paragraph().text = "Redundancy: stock 20% spare drivetrains + batteries to keep fleet field-ready."

# Slide 10 Timeline & Milestones
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "30-60-90 Day Milestones"
body = slide.shapes.placeholders[1].text_frame
body.text = "0–30 days: hardware build + 5 successful outdoor pacing runs"
body.add_paragraph().text = "31–60 days: reliability hardening, pilot workflow, telemetry dashboards"
body.add_paragraph().text = "61–90 days: investor-ready pilot package (demo script, KPI dashboard, video proof)"

# Slide 11 Collaboration invites
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Collaboration Opportunities"
cols = ["Investors", "Developers", "Hardware Team"]
points = [
    ["Fund hardware + pilot ops", "Help shape go-to-market", "Access to progress reporting"],
    ["Extend cloud module surfaces", "Integrate partner APIs", "Enhance AI coaching experiences"],
    ["Finalize wiring & enclosure", "Optimize power + drivetrain", "Stress-test outdoor durability"],
]
for i, (col, bullet_list) in enumerate(zip(cols, points)):
    box = slide.shapes.add_textbox(Inches(0.5 + i * 3), Inches(1.7), Inches(2.6), Inches(3.5))
    tf = box.text_frame
    tf.text = col
    tf.paragraphs[0].font.bold = True
    for b in bullet_list:
        p = tf.add_paragraph()
        p.text = b
        p.level = 1

# Slide 12 Financial Projections
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Financial Projections"
body = slide.shapes.placeholders[1].text_frame
body.text = "Year 1 pilot: 20 units @ $300/mo subscription → $72k ARR"
body.add_paragraph().text = "Year 2 regional rollout: 150 units, $450/mo avg → $810k ARR"
body.add_paragraph().text = "Hardware COGS per unit: $325 prototype → $210 scaled manufacturing"
body.add_paragraph().text = "Gross margin target: 65% with hybrid software upsells (AI coach, club dashboards)"

# Slide 13 Competitive Landscape
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Competitive Landscape"
body = slide.shapes.placeholders[1].text_frame
body.text = "No current product combines physical pacing robot + hybrid module software."
for item in [
    "Treadmill/connected fitness: indoor only, lacks outdoor alerting",
    "Wearable coaching apps: software-only, cannot physically pace or clear path",
    "Delivery robots (Starship, Serve): not optimized for runners, no voice/alert focus",
]:
    body.add_paragraph().text = item
body.add_paragraph().text = "Moat: voice-authenticated control + safety alerts + OTA modules for partners."

# Slide 14 Team Bios
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Core Team"
body = slide.shapes.placeholders[1].text_frame
body.text = "Eddie Lee — Robotics + mobile lead (ex-Garmin, built sports telemetry stacks)."
body.add_paragraph().text = "Maya Patel — Safety/compliance (former race director, EMT experience)."
body.add_paragraph().text = "Jamal Ortiz — AI/voice systems (deployed multi-modal copilots at enterprise scale)."
body.add_paragraph().text = "Advisors: collegiate track coach network + hardware manufacturing partner in Shenzhen."

# Slide 15 Funding Ask
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Funding & Support Ask"
body = slide.shapes.placeholders[1].text_frame
body.text = "Seed support requested to:"
for item in [
    "Procure hardware fleet + test logistics ($5k)",
    "Stand up hybrid content control plane + analytics ($8k)",
    "Cover pilot operations and safety compliance ($7k)",
]:
    body.add_paragraph().text = item
body.add_paragraph().text = "Optional strategic partnership: running clubs, race organizers, safety tech vendors."

# Slide 16 Call to Action
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Next Steps"
body = slide.shapes.placeholders[1].text_frame
body.text = "1. Approve hardware procurement + pilot budget"
body.add_paragraph().text = "2. Schedule hybrid module workshop (investor + developer views)"
body.add_paragraph().text = "3. Confirm pilot routes and stakeholder contacts"
body.add_paragraph().text = "4. Target demo day in 90 days with full field footage"

prs.save('c:/src/runner_companion_app/docs/RunnerCompanion_Pitch.pptx')
